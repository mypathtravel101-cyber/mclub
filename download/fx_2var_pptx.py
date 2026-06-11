
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json

with open('/home/z/my-project/download/ppt_charts/stress_2var_data.json','r') as f:
    D = json.load(f)

CH = '/home/z/my-project/download/ppt_charts'
GOLD=RGBColor(0xC9,0xA8,0x4C); DARK=RGBColor(0x1A,0x1A,0x1A); BLUE=RGBColor(0x1E,0x3A,0x5F)
WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xF5,0xF2,0xE8); RED=RGBColor(0xE7,0x4C,0x3C)
GREEN=RGBColor(0x2E,0xCC,0x71); GREY=RGBColor(0x95,0xA5,0xA6); ORANGE=RGBColor(0xF3,0x9C,0x12)
PURPLE=RGBColor(0x8E,0x44,0xAD)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s,c=None):
    if c is None: c=DARK
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def bar(s,t=None):
    if t is None: t=Inches(1.2)
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),t,Inches(13.333),Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb=GOLD; sh.line.fill.background()
def tb(s,l,t,w,h,txt,sz=14,b=False,c=None,a=None):
    if c is None: c=WHITE
    if a is None: a=PP_ALIGN.LEFT
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.bold=b
    p.font.color.rgb=c; p.font.name='Sarasa Mono SC'; p.alignment=a; return bx
def mkbox(s,l,t,w,h,fc,bc=None,bw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if bc:
        sh.line.color.rgb=bc
        sh.line.width = bw if bw else Pt(2)
    else:
        sh.line.fill.background()
    return sh
def ft(s):
    tb(s,Inches(0.5),Inches(6.9),Inches(12),Inches(0.4),
       'PZC Group \u767e\u76db\u5927\u901a \u2014 MCLUB FX Risk Modeling Service | \u96d9\u8b8a\u91cf\u58d3\u529b\u6e2c\u8a66\u5831\u544a',sz=9,c=GREY,a=PP_ALIGN.CENTER)

# SLIDE 1: COVER
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); bar(sl,Inches(2.6))
tb(sl,Inches(1),Inches(0.8),Inches(11),Inches(0.5),'MCLUB FX RISK MODELING SERVICE',sz=16,c=GOLD,a=PP_ALIGN.CENTER)
tb(sl,Inches(1),Inches(2.9),Inches(11),Inches(1.5),'\u65e5\u672c\u7269\u696d\u6295\u8cc7 \u2014 \u96d9\u8b8a\u91cf\u58d3\u529b\u6e2c\u8a66',sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(sl,Inches(1),Inches(4.2),Inches(11),Inches(0.5),'2-Variable Stress Test: FX Rate x Property Value',sz=16,c=GREY,a=PP_ALIGN.CENTER)

cv = "CONSTANT \u56fa\u5b9a\u6578\u64da\n\n\u7269\u696d: HKD 4,000,000 (JPY 78M)\n\u6309\u63d0: 40% LTV @ 3% p.a.\n\u671f\u9650: 15\u5e74 (45\u219260\u6b72)\n\u6708\u4f9b: JPY 215,461\n\u5e74\u79df\u91d1: JPY 4,680,000 (6%)\n\u5ba2\u6236\u6295\u5165: HKD 2,400,000"
vr = "VARIABLE \u8b8a\u91cf\u6578\u64da\n\nVariable 1: FX Rate (JPY/HKD)\n  12\u500b\u60c5\u666f: T1/T2/T3 + \u6b77\u53f2\n  \u7bc4\u570d: JPY +25% ~ -31%\n\nVariable 2: Property Value\n  7\u500b\u60c5\u666f: \u5d29\u76e4~\u589e\u9577+80%\n  \u7bc4\u570d: 15\u5e74 -30% ~ +80%"

mkbox(sl,Inches(1.5),Inches(5.0),Inches(4.5),Inches(1.7),RGBColor(0x2C,0x3E,0x50),BLUE)
tb(sl,Inches(1.7),Inches(5.05),Inches(4.1),Inches(1.6),cv,sz=10,c=LIGHT)
mkbox(sl,Inches(7.3),Inches(5.0),Inches(4.5),Inches(1.7),RGBColor(0x2C,0x3E,0x50),ORANGE)
tb(sl,Inches(7.5),Inches(5.05),Inches(4.1),Inches(1.6),vr,sz=10,c=LIGHT)
ft(sl)

# SLIDE 2: 2D HEATMAP
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),'\u96d9\u8b8a\u91cf\u58d3\u529b\u6e2c\u8a66 \u2014 84\u500b\u60c5\u666f\u5e74\u5316\u56de\u5831\u77e9\u9663',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))
sl.shapes.add_picture(f'{CH}/stress_2var_heatmap.png',Inches(0.2),Inches(1.0),Inches(12.9),Inches(5.6))
mkbox(sl,Inches(0.5),Inches(6.65),Inches(3.8),Inches(0.5),RGBColor(0x1A,0x3C,0x34))
tb(sl,Inches(0.6),Inches(6.68),Inches(3.6),Inches(0.4),f'\u6700\u5dee\u60c5\u666f: {D["worst"]["annualized"]}% \u5e74\u5316',sz=11,b=True,c=RED,a=PP_ALIGN.CENTER)
mkbox(sl,Inches(4.7),Inches(6.65),Inches(3.8),Inches(0.5),RGBColor(0x2C,0x2E,0x30))
tb(sl,Inches(4.8),Inches(6.68),Inches(3.6),Inches(0.4),'\u866b\u672c\u60c5\u666f: 0/84 (0%)',sz=11,b=True,c=GREEN,a=PP_ALIGN.CENTER)
mkbox(sl,Inches(8.9),Inches(6.65),Inches(3.9),Inches(0.5),RGBColor(0x1A,0x3C,0x34))
tb(sl,Inches(9.0),Inches(6.68),Inches(3.7),Inches(0.4),'\u57fa\u6e96\u60c5\u666f: 5.83% \u5e74\u5316',sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
ft(sl)

# SLIDE 3: SENSITIVITY LINES
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),'\u4e0d\u540c\u623f\u50f9\u60c5\u666f\u4e0b\u7684\u532f\u7387\u654f\u611f\u5ea6\u66f2\u7dda',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))
sl.shapes.add_picture(f'{CH}/stress_2var_lines.png',Inches(0.3),Inches(1.0),Inches(12.7),Inches(5.6))
ft(sl)

# SLIDE 4: CONTOUR / RISK MAP
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),'\u76c8\u8667\u5e73\u8861\u98a8\u96aa\u5730\u5716 \u2014 FX Rate x Property Value',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))
sl.shapes.add_picture(f'{CH}/stress_2var_contour.png',Inches(0.3),Inches(1.0),Inches(12.7),Inches(5.6))
ft(sl)

# SLIDE 5: FULL RESULTS TABLE
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),'\u58d3\u529b\u6e2c\u8a66\u5b8c\u6574\u7d50\u679c\u8868',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

fx_labels = [f['label'] for f in D['variable1_fx']]
prop_labels = [p['label'] for p in D['variable2_property']]
mat = D['matrix_annualized']

rows = len(prop_labels) + 1
cols = len(fx_labels) + 1
tbl = sl.shapes.add_table(rows, cols, Inches(0.15), Inches(1.05), Inches(13.0), Inches(5.5)).table

tbl.columns[0].width = Inches(1.4)
cw = (13.0-1.4)/(cols-1)
for i in range(1, cols):
    tbl.columns[i].width = Inches(cw)

tbl.cell(0,0).text = '\u623f\u50f9\\\u532f\u7387'
for j, fl in enumerate(fx_labels):
    tbl.cell(0,j+1).text = fl
for j in range(cols):
    cell=tbl.cell(0,j)
    for p in cell.text_frame.paragraphs:
        p.font.size=Pt(8.5); p.font.bold=True; p.font.color.rgb=GOLD
        p.font.name='Sarasa Mono SC'; p.alignment=PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0x2C,0x3E,0x50)

for ri, (pl, row_vals) in enumerate(zip(prop_labels, mat)):
    tbl.cell(ri+1,0).text = pl
    cell0=tbl.cell(ri+1,0)
    for p in cell0.text_frame.paragraphs:
        p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=WHITE
        p.font.name='Sarasa Mono SC'
    cell0.fill.solid(); cell0.fill.fore_color.rgb=RGBColor(0x25,0x25,0x30)
    
    for ci, v in enumerate(row_vals):
        cell=tbl.cell(ri+1,ci+1); cell.text=f'{v:.2f}%'
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(8.5); p.font.name='Sarasa Mono SC'; p.alignment=PP_ALIGN.CENTER
            p.font.bold = (ci == 5 and ri == 2)
            if v >= 8: p.font.color.rgb = GREEN
            elif v >= 6: p.font.color.rgb = RGBColor(0x82,0xE0,0xAA)
            elif v >= 5: p.font.color.rgb = GOLD
            elif v >= 3: p.font.color.rgb = ORANGE
            else: p.font.color.rgb = RED
        cell.fill.solid()
        if ci == 5 and ri == 2:
            cell.fill.fore_color.rgb = RGBColor(0x3A,0x35,0x20)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x20,0x22,0x28)
ft(sl)

# SLIDE 6: CONCLUSIONS
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),'\u58d3\u529b\u6e2c\u8a66\u7d50\u8ad6\u8207\u98a8\u96aa\u8a55\u4f30',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

findings = [
    ("\u6e2c\u8a66\u7bc4\u570d", "12\u500b\u532f\u7387\u60c5\u666f x 7\u500b\u623f\u50f9\u60c5\u666f = 84\u500b\u58d3\u529b\u6e2c\u8a66\u60c5\u666f\n\u8986\u84cbMCLUB\u4e09\u5c64\u67b6\u69cb\uff08T1/T2/T3\uff09+ 2022\u6b77\u53f2\u53c3\u7167", BLUE),
    ("\u96f6\u866b\u672c\u6982\u7387", "\u5168\u90e884\u500b\u60c5\u666f\u5747\u9304\u5f97\u6b63\u56de\u5831\uff0c\u6700\u5dee\u5e74\u5316\u4ecd\u90542.20%\n\u5373\u4f7f\u623f\u50f9\u5d29\u76e430% + \u65e5\u5703\u8cb6\u503c31%\u540c\u6642\u767c\u751f\uff0c\u6295\u8cc7\u4ecd\u53ef\u4fdd\u672c", GREEN),
    ("\u57fa\u6e96\u60c5\u666f", "\u623f\u50f9\u6301\u5e73 + \u73fe\u532f\u7387 = \u5e74\u53165.83%\uff08\u7d14\u79df\u91d1\uff09\n\u82e5\u623f\u50f9\u5e74\u589e3% + \u73fe\u532f\u7387 = \u5e74\u53168.21%\uff08\u542b\u8cc7\u672c\u589e\u503c\uff09", GOLD),
    ("\u81ea\u7136\u5c0d\u6c96\u6548\u61c9", "\u65e5\u5703\u6309\u63d0+\u65e5\u5703\u79df\u91d1\u5f62\u6210\u5e63\u7a2e\u5339\u914d\u5c0d\u6c96\n\u532f\u7387\u98a8\u96aa\u50c5\u5f71\u97ff\u6de8\u79df\u91d1\u53ca\u5957\u73fe\u50f9\u503c\uff0c\u975e\u5168\u90e8\u7269\u696d\u50f9\u503c", BLUE),
    ("\u69d3\u6975\u5b89\u5168\u908a\u969b", "40% LTV\u63d0\u4f9b\u5145\u8db3\u7de9\u885d\uff0c\u7269\u696d\u9700\u8cb6\u503c\u8d8560%\u624d\u89f8\u53ca\u8cb8\u6b3e\u5b89\u5168\u7dda\n\u6309\u63d0\u672c\u91d1\u511f\u9084\u6301\u7e8c\u964d\u4f4e\u8cb8\u6b3e\u9918\u984d\uff0c15\u5e74\u5f8c\u5b8c\u5168\u6e05\u511f", RGBColor(0x34,0x98,0xDB)),
    ("\u96d9\u91cd\u98a8\u96aa\u6700\u5dee\u60c5\u666f", "\u623f\u50f9-30% + JPY\u8cb631%\uff082022\u5371\u6a5f\uff09\n\u5e74\u53162.20%\uff0c15\u5e74\u7e3dROI +38.8%\uff0c\u4ecd\u5927\u5e45\u8dd1\u8d0f\u901a\u8132", RED),
]

y = Inches(1.15)
for title, desc, color in findings:
    mkbox(sl, Inches(0.5), y, Inches(7.5), Inches(0.9), RGBColor(0x2C,0x3E,0x50), color)
    tb(sl, Inches(0.7), y+Inches(0.05), Inches(2.0), Inches(0.35), title, sz=12, b=True, c=color)
    tb(sl, Inches(0.7), y+Inches(0.35), Inches(7.0), Inches(0.55), desc, sz=9.5, c=LIGHT)
    y += Inches(0.95)

# Right: Summary card
mkbox(sl,Inches(8.3),Inches(1.15),Inches(4.5),Inches(5.5),RGBColor(0x2C,0x3E,0x50),GOLD)
tb(sl,Inches(8.5),Inches(1.25),Inches(4.1),Inches(0.4),'\u98a8\u96aa\u8a55\u7d1a\u7e3d\u7d50',sz=16,b=True,c=GOLD,a=PP_ALIGN.CENTER)

grades = [
    ("\u6574\u9ad4\u98a8\u96aa\u7b49\u7d1a", "LOW-MODERATE\n\u4e2d\u7b49\u504f\u4f4e", GREEN),
    ("FX\u98a8\u96aa", "MITIGATED\n\u5df2\u88ab\u81ea\u7136\u5c0d\u6c96\u7de9\u89e3", RGBColor(0x82,0xE0,0xAA)),
    ("\u623f\u50f9\u98a8\u96aa", "LOW-MODERATE\n15\u5e74\u9577\u671f\u6301\u6709\u5206\u6563\u98a8\u96aa", GOLD),
    ("\u5229\u7387\u98a8\u96aa", "LOW\n\u65e5\u672c3%\u8655\u6b77\u53f2\u4f4e\u4f4d", GREEN),
    ("\u6d41\u52d5\u6027\u98a8\u96aa", "MODERATE\n\u65e5\u672c\u7269\u696d\u6d41\u52d5\u6027\u504f\u4f4e", ORANGE),
    ("\u5efa\u8b70", "\u9069\u7528MCLUB\u9032\u968e\u7248\n(Advanced)\u670d\u52d9\u5c64\u7d1a\nHK$128,000", PURPLE),
]

gy = Inches(1.75)
for label, val, color in grades:
    tb(sl, Inches(8.5), gy, Inches(4.1), Inches(0.3), label, sz=10, c=GREY)
    tb(sl, Inches(8.5), gy+Inches(0.28), Inches(4.1), Inches(0.5), val, sz=13, b=True, c=color)
    gy += Inches(0.82)

ft(sl)

out = '/home/z/my-project/download/PZC_FX_2Var_StressTest_JPY_Property.pptx'
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
