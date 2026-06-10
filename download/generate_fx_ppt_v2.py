#!/usr/bin/env python3
"""
Generate an EDITABLE 10-slide PPTX for PZC Group FX Risk Modelling.
All text uses native PowerPoint text placeholders — fully editable in PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette (IG-1 Ink Gold) ──
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF5, 0xF2, 0xE8)
GRAY      = RGBColor(0x99, 0x99, 0x99)
MED_GRAY  = RGBColor(0x66, 0x66, 0x66)
BLACK     = RGBColor(0x00, 0x00, 0x00)

OUT = "/home/z/my-project/download/PZC_Group_FX_Risk_Modelling_可编辑版.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── helpers ──────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    """Add a plain textbox — always fully editable."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=BLACK, alignment=PP_ALIGN.LEFT,
                          font_name="Microsoft YaHei", line_spacing=1.5,
                          bold_first=False):
    """Add textbox with multiple paragraphs (each line = one paragraph)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(font_size * 0.4)
        pf = p._pPr
        if pf is None:
            from lxml import etree
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            pf = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        lnSpc = etree.SubElement(pf, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox

def add_shape_bg(slide, left, top, width, height, fill_color):
    """Add a coloured rectangle background shape."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_gold_bar(slide, top=Inches(1.15)):
    """Add a thin gold accent line under the slide header."""
    add_shape_bg(slide, Inches(0.8), top, Inches(2), Pt(3), GOLD)

def slide_header(slide, title_text, subtitle_text=""):
    """Standard slide header with gold bar."""
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                title_text, font_size=28, bold=True, color=DARK,
                font_name="SimHei")
    add_gold_bar(slide)
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
                    subtitle_text, font_size=14, color=MED_GRAY)

def add_arrow_right(slide, left, top, width, height, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_chevron(slide, left, top, width, height, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_circle(slide, left, top, size, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(sl, 0, 0, W, H, DARK)
add_shape_bg(sl, 0, H - Inches(1.2), W, Inches(1.2), GOLD)

add_textbox(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
            "PZC Group", font_size=44, bold=True, color=GOLD,
            font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(2.6), Inches(10), Inches(1.0),
            "FX Risk Modelling", font_size=36, bold=True, color=WHITE,
            font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
            "外匯風險建模服務簡介", font_size=24, color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(4.4), Inches(10), Inches(0.5),
            "全面外匯風險管理解決方案  |  服務家族客戶外匯資產保值",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
            "百盛大通  |  专业外汇风险管理服务", font_size=14, color=DARK,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — OVERVIEW (4 services at a glance)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "四大核心服務概覽", "FX Risk Modelling — 四項專業外匯風險管理服務")

cards = [
    ("01", "壓力測試\nStress Testing", "模擬極端市場情況，\n測試外匯資產組合損失風險"),
    ("02", "多幣種儀表板\nMulti-Currency Dashboard", "實時多幣種分析平台，\n一站式監控全球匯率走勢"),
    ("03", "對沖策略\nHedging Strategies", "個性化對沖方案，\n最大限度降低匯率波動侵蝕"),
    ("04", "風險估值\nValue at Risk (VaR)", "統計模型量化風險，\n提供簡潔明瞭的風險指標"),
]
x_start = Inches(0.6)
card_w = Inches(2.9)
gap = Inches(0.25)
card_h = Inches(4.2)
y_top = Inches(2.2)

for i, (num, title, desc) in enumerate(cards):
    x = x_start + i * (card_w + gap)
    add_rounded_rect(sl, x, y_top, card_w, card_h, LIGHT_BG)
    add_circle(sl, x + Inches(0.2), y_top + Inches(0.2), Inches(0.6), GOLD)
    add_textbox(sl, x + Inches(0.2), y_top + Inches(0.25), Inches(0.6), Inches(0.55),
                num, font_size=18, bold=True, color=DARK,
                alignment=PP_ALIGN.CENTER, font_name="Calibri")
    add_multiline_textbox(sl, x + Inches(0.2), y_top + Inches(1.0), card_w - Inches(0.4), Inches(0.9),
                          title.split("\n"), font_size=16, bold_first=True, color=DARK,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    add_multiline_textbox(sl, x + Inches(0.2), y_top + Inches(2.1), card_w - Inches(0.4), Inches(1.8),
                          desc.split("\n"), font_size=13, color=MED_GRAY,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — STRESS TESTING (detail)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務一：壓力測試 Stress Testing")

# Left content
add_multiline_textbox(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    ["模擬極端市場情況（如幣值急跌 10%-20%、利率劇變、",
     "地緣政治事件），測試客戶的外匯資產組合在各種情境下",
     "的損失風險。",
     "",
     "幫助客戶提前知悉最壞情況下的潛在影響，並制定應對方案。",
     "",
     "服務對象：有大量外匯負債的家族客戶",
     "",
     "交付成果：壓力測試報告（包含情境分析、損失估算、建議）"],
    font_size=16, color=DARK, line_spacing=1.6)

# Right — process diagram (3 rounded boxes with arrows)
steps = ["定義壓力情境\n(極端市場事件)", "模擬損失計算\n(資產組合影響)", "制定應對方案\n(報告與建議)"]
box_w = Inches(2.2)
box_h = Inches(1.2)
y_st = Inches(2.8)
x_base = Inches(7.2)

for j, st in enumerate(steps):
    x = x_base + j * Inches(2.6)
    add_rounded_rect(sl, x, y_st, box_w, box_h, GOLD)
    add_multiline_textbox(sl, x + Inches(0.1), y_st + Inches(0.15), box_w - Inches(0.2), box_h - Inches(0.3),
                          st.split("\n"), font_size=13, bold_first=True, color=DARK,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    if j < 2:
        add_arrow_right(sl, x + box_w + Inches(0.1), y_st + Inches(0.35),
                        Inches(0.35), Inches(0.5), GOLD)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — MULTI-CURRENCY DASHBOARD (detail)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務二：多幣種儀表板 Multi-Currency Dashboard")

add_multiline_textbox(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    ["實時多幣種分析平台，整合威利指定價、幣值走勢、",
     "外匯敞口數據。",
     "",
     "客戶可一網盯控全球幣種走勢和自己的外匯風險敏感度，",
     "即時追蹤幣值波動對負債的實際影響。",
     "",
     "服務對象：會經常要進行多幣種交易的客戶",
     "",
     "交付成果：實時儀表板（幣值走勢、敞口分析、風險指標）"],
    font_size=16, color=DARK, line_spacing=1.6)

# Right — dashboard mock-up diagram
add_rounded_rect(sl, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), LIGHT_BG)
add_textbox(sl, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.5),
            "Dashboard 示意圖", font_size=14, bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)

# Mini chart bars (simulate bar chart using rectangles)
bar_data = [("USD", 0.7), ("EUR", 0.55), ("GBP", 0.45), ("JPY", 0.65), ("CNY", 0.8)]
bar_w = Inches(0.6)
bar_base_y = Inches(5.5)
max_h = Inches(2.0)
for k, (label, val) in enumerate(bar_data):
    bx = Inches(7.5) + k * Inches(0.9)
    bh = max_h * val
    add_shape_bg(sl, bx, bar_base_y - bh, bar_w, bh, GOLD)
    add_textbox(sl, bx - Inches(0.05), bar_base_y + Inches(0.05), bar_w + Inches(0.1), Inches(0.3),
                label, font_size=11, color=DARK, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# ════════════════════════════════════════════════════════════
# SLIDE 5 — HEDGING STRATEGIES (detail)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務三：對沖策略 Hedging Strategies")

add_multiline_textbox(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    ["根據客戶的外匯敞口和風險實際度，設計個性化的對沖方案。",
     "",
     "包括遠期合約、期權策略、天然對沖等多種工具組合，",
     "最大限度降低匯率波動對家族資產的侵蝕。",
     "",
     "服務對象：希望主動管理匯率風險的客戶",
     "",
     "交付成果：對沖方案建議書（包含成本效益分析）"],
    font_size=16, color=DARK, line_spacing=1.6)

# Right — hedging tools diagram
tools = ["遠期合約\nForward", "期權策略\nOptions", "天然對沖\nNatural Hedge", "互換合約\nSwaps"]
y_tools = Inches(2.8)
for j, tool in enumerate(tools):
    x = Inches(7.2) + j * Inches(1.5)
    add_rounded_rect(sl, x, y_tools, Inches(1.3), Inches(1.0), GOLD)
    add_multiline_textbox(sl, x + Inches(0.05), y_tools + Inches(0.1), Inches(1.2), Inches(0.8),
                          tool.split("\n"), font_size=12, bold_first=True, color=DARK,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.2)

# Arrow pointing down to result
add_arrow_right(sl, Inches(9.5), Inches(4.2), Inches(1.0), Inches(0.6), GOLD)
# Actually let's use a down arrow
from pptx.enum.shapes import MSO_SHAPE
down_arrow = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.7), Inches(4.2), Inches(0.6), Inches(0.6))
down_arrow.fill.solid()
down_arrow.fill.fore_color.rgb = GOLD
down_arrow.line.fill.background()

add_rounded_rect(sl, Inches(7.5), Inches(5.0), Inches(5.0), Inches(1.0), LIGHT_BG)
add_textbox(sl, Inches(7.7), Inches(5.15), Inches(4.6), Inches(0.7),
            "組合運用 → 最小化匯率風險，最大化資產保護",
            font_size=14, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — VALUE AT RISK (detail)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務四：風險估值 Value at Risk (VaR)")

add_multiline_textbox(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    ["利用統計模型計算客戶外匯組合在特定置信水平下的",
     "最大預期損失。",
     "",
     "例如「95% VaR 為 HK$500萬」意味在正常市場條件下，",
     "有 95% 機會單日損失不會超過 HK$500萬。",
     "",
     "服務對象：需要定期報告外匯風險敞口的客戶",
     "",
     "交付成果：VaR 報告（包含置信度、風險值、走勢圖）"],
    font_size=16, color=DARK, line_spacing=1.6)

# Right — VaR bell curve diagram (simplified with shapes)
# Draw a normal distribution shape using an arc
add_rounded_rect(sl, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), LIGHT_BG)
add_textbox(sl, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.5),
            "VaR 概念示意圖", font_size=14, bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)

# Bell curve representation
# X-axis
add_shape_bg(sl, Inches(7.5), Inches(5.5), Inches(4.5), Pt(2), DARK)
# 95% region highlight
add_shape_bg(sl, Inches(8.0), Inches(3.5), Inches(3.5), Inches(2.0),
             RGBColor(0xC9, 0xA8, 0x4C))  # Gold area

add_textbox(sl, Inches(8.5), Inches(3.7), Inches(2.5), Inches(1.2),
            "95% 置信區間", font_size=16, bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(8.5), Inches(4.3), Inches(2.5), Inches(0.5),
            "VaR = HK$500萬", font_size=14, color=DARK,
            alignment=PP_ALIGN.CENTER)

# Labels
add_textbox(sl, Inches(7.5), Inches(5.6), Inches(1.0), Inches(0.3),
            "-2σ", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(sl, Inches(9.5), Inches(5.6), Inches(1.0), Inches(0.3),
            "均值", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(11.0), Inches(5.6), Inches(1.0), Inches(0.3),
            "+2σ", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# ════════════════════════════════════════════════════════════
# SLIDE 7 — SERVICE COMPARISON TABLE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務對比總覽", "四大核心服務一覽表")

# Table
rows_data = [
    ["服務項目", "核心功能", "服務對象", "交付成果"],
    ["壓力測試\nStress Testing", "模擬極端市場情況\n測試損失風險", "有大量外匯負債\n的家族客戶", "壓力測試報告\n(情境分析/損失估算/建議)"],
    ["多幣種儀表板\nDashboard", "實時多幣種分析\n一網監控匯率走勢", "經常進行多幣種\n交易的客戶", "實時儀表板\n(走勢/敞口/風險指標)"],
    ["對沖策略\nHedging", "個性化對沖方案\n多工具組合", "希望主動管理\n匯率風險的客戶", "對沖方案建議書\n(成本效益分析)"],
    ["風險估值\nVaR", "統計模型量化風險\n簡潔風險指標", "需定期報告外匯\n風險敞口的客戶", "VaR 報告\n(置信度/風險值/走勢圖)"],
]

tbl = sl.shapes.add_table(len(rows_data), 4, Inches(0.6), Inches(2.0),
                           Inches(12), Inches(5.0)).table
tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(3.7)

for ri, row in enumerate(rows_data):
    for ci, cell_text in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Microsoft YaHei"
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = DARK
                p.font.size = Pt(13)
            else:
                p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
        # Header row styling
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD
        elif ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# ════════════════════════════════════════════════════════════
# SLIDE 8 — WORKFLOW PROCESS DIAGRAM
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "服務流程圖", "PZC Group FX Risk Modelling 工作流程")

steps_flow = [
    "初步評估\n風險敞口分析",
    "定制方案\n選擇合適服務組合",
    "模型建設\n數據整合與計算",
    "報告交付\n專業分析報告",
    "持續跟進\n定期更新與調整",
]
step_w = Inches(2.0)
step_h = Inches(1.4)
y_flow = Inches(3.0)
x_flow_start = Inches(0.5)

for i, step in enumerate(steps_flow):
    x = x_flow_start + i * Inches(2.5)
    # Circle with step number
    add_circle(sl, x + Inches(0.65), Inches(2.2), Inches(0.7), GOLD)
    add_textbox(sl, x + Inches(0.65), Inches(2.25), Inches(0.7), Inches(0.6),
                str(i + 1), font_size=22, bold=True, color=DARK,
                alignment=PP_ALIGN.CENTER, font_name="Calibri")
    # Step box
    add_rounded_rect(sl, x, y_flow, step_w, step_h, LIGHT_BG)
    add_multiline_textbox(sl, x + Inches(0.1), y_flow + Inches(0.15), step_w - Inches(0.2), step_h - Inches(0.3),
                          step.split("\n"), font_size=14, bold_first=True, color=DARK,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    # Arrow between steps
    if i < len(steps_flow) - 1:
        add_arrow_right(sl, x + step_w + Inches(0.1), y_flow + Inches(0.4),
                        Inches(0.35), Inches(0.5), GOLD)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — WHY PZC GROUP
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "為什麼選擇 PZC Group？", "FX Risk Modelling 核心優勢")

advantages = [
    ("專業團隊", "擁有資深外匯分析師團隊，\n具備豐富市場經驗與專業知識"),
    ("定制服務", "根據每位客戶的獨特需求，\n量身打造風險管理方案"),
    ("科技驅動", "先進的數據分析平台，\n精準計算與即時監控"),
    ("全面覆蓋", "四大服務全方位覆蓋外匯\n風險管理的各個關鍵環節"),
]

adv_w = Inches(2.7)
adv_h = Inches(3.0)
y_adv = Inches(2.2)
x_adv_start = Inches(0.6)
adv_gap = Inches(0.35)

for i, (title, desc) in enumerate(advantages):
    x = x_adv_start + i * (adv_w + adv_gap)
    add_rounded_rect(sl, x, y_adv, adv_w, adv_h, LIGHT_BG)
    # Gold top bar
    add_shape_bg(sl, x, y_adv, adv_w, Inches(0.08), GOLD)
    add_textbox(sl, x + Inches(0.2), y_adv + Inches(0.3), adv_w - Inches(0.4), Inches(0.6),
                title, font_size=20, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(sl, x + Inches(0.2), y_adv + Inches(1.2), adv_w - Inches(0.4), Inches(1.5),
                          desc.split("\n"), font_size=14, color=MED_GRAY,
                          alignment=PP_ALIGN.CENTER, line_spacing=1.6)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — CONTACT / CLOSING
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(sl, 0, 0, W, H, DARK)
add_shape_bg(sl, 0, H - Inches(1.2), W, Inches(1.2), GOLD)

add_textbox(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
            "PZC Group", font_size=44, bold=True, color=GOLD,
            font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8),
            "FX Risk Modelling", font_size=32, bold=True, color=WHITE,
            font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
            "外匯風險建模服務", font_size=20, color=WHITE,
            alignment=PP_ALIGN.CENTER)

# Contact info (editable text)
add_multiline_textbox(sl, Inches(2.5), Inches(4.4), Inches(8), Inches(1.8),
    ["如需了解更多詳情，歡迎聯繫我們",
     "",
     "（請在此處填入聯繫方式）"],
    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_textbox(sl, Inches(1.5), Inches(5.9), Inches(10), Inches(0.4),
            "百盛大通  |  专业外汇风险管理服务", font_size=14, color=DARK,
            alignment=PP_ALIGN.CENTER)

# ── SAVE ───────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved to {OUT}")
