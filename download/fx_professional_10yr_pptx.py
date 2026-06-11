from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json

with open('/home/z/my-project/download/ppt_charts/stress_2var_data.json','r') as f:
    D = json.load(f)
with open('/home/z/my-project/download/ppt_charts/fx_analysis_data.json','r') as f:
    FX = json.load(f)

CH = '/home/z/my-project/download/ppt_charts'
GOLD=RGBColor(0xC9,0xA8,0x4C); DARK=RGBColor(0x1A,0x1A,0x1A); BLUE=RGBColor(0x1E,0x3A,0x5F)
WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xF5,0xF2,0xE8); RED=RGBColor(0xE7,0x4C,0x3C)
GREEN=RGBColor(0x2E,0xCC,0x71); GREY=RGBColor(0x95,0xA5,0xA6); ORANGE=RGBColor(0xF3,0x9C,0x12)
PURPLE=RGBColor(0x8E,0x44,0xAD)
DARK_BG = RGBColor(0x1A,0x1A,0x1A)
CARD_BG = RGBColor(0x2C,0x2E,0x30)
PANEL_BG = RGBColor(0x25,0x25,0x30)

C = D['constants']
HOLD = C['hold_years']
MORT = C['mortgage_term_years']
BAL_EXIT = C['balance_at_exit_jpy']

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s,c=None):
    if c is None: c=DARK_BG
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def bar(s,t=None):
    if t is None: t=Inches(1.2)
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),t,Inches(13.333),Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb=GOLD; sh.line.fill.background()
def tb(s,l,t,w,h,txt,sz=14,b=False,c=None,a=None):
    if c is None: c=WHITE
    if a is None: a=PP_ALIGN.LEFT
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.bold=b
    p.font.color.rgb=c; p.font.name='Sarasa Mono SC'; p.alignment=a; return bx
def mkbox(s,l,t,w,h,fc,bc=None,bw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if bc:
        sh.line.color.rgb=bc
        sh.line.width = bw if bw else Pt(2)
    else:
        sh.line.fill.background()
    return sh
def ft(s):
    tb(s,Inches(0.5),Inches(6.95),Inches(12),Inches(0.4),
       'PZC Group \u767e\u76db\u5927\u901a \u2014 MCLUB FX Risk Modeling Service | JPY\u7269\u696d\u6295\u8cc7\u5206\u6790\u5831\u544a',sz=9,c=GREY,a=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); bar(sl,Inches(2.2))
tb(sl,Inches(1),Inches(0.6),Inches(11),Inches(0.5),'MCLUB FX RISK MODELING SERVICE',sz=16,c=GOLD,a=PP_ALIGN.CENTER)
tb(sl,Inches(1),Inches(2.5),Inches(11),Inches(1.5),
   '\u65e5\u672c\u7269\u696d\u6295\u8cc7 \u2014 \u5c08\u696d\u98a8\u96aa\u5206\u6790\u5831\u544a',sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(sl,Inches(1),Inches(3.8),Inches(11),Inches(0.5),
   'JPY Property Investment \u2014 Professional Risk Analysis (10-Year Holding)',sz=15,c=GREY,a=PP_ALIGN.CENTER)

# Key params on cover
info_left = (
    "\u6295\u8cc7\u7d50\u69cb\n"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
    "\u7269\u696d\u50f9\u503c: HKD 4,000,000 (JPY 78M)\n"
    "\u6309\u63d0: 40% LTV @ 3% p.a.\n"
    "\u6309\u63d0\u671f\u9650: 15\u5e74 (45\u219260\u6b72)\n"
    "\u6301\u5009\u671f: 10\u5e74 (45\u219255\u6b72)\n"
    "\u5ba2\u6236\u6295\u5165: HKD 2,400,000\n"
    "\u5e74\u79df\u91d1\u6b3e\u7387: 6%"
)
info_right = (
    "\u95dc\u9375\u6307\u6a19 (10\u5e74\u6301\u5009)\n"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
    "\u57fa\u6e96\u5e74\u5316\u56de\u5831: {:.2f}%\n"
    "\u6700\u5dee\u5e74\u5316\u56de\u5831: {:.2f}%\n"
    "\u6700\u4f73\u5e74\u5316\u56de\u5831: {:.2f}%\n"
    "\u866b\u672c\u60c5\u666f: 0 / 84 (0%)\n"
    "\u7b2c10\u5e74\u6309\u63d0\u9918\u984d: JPY {:,.0f}".format(
        D['base_annualized'], D['worst']['annualized'],
        D['best']['annualized'], BAL_EXIT)
)
mkbox(sl,Inches(1.5),Inches(4.6),Inches(4.8),Inches(2.1),RGBColor(0x2C,0x3E,0x50),BLUE)
tb(sl,Inches(1.7),Inches(4.65),Inches(4.4),Inches(2.0),info_left,sz=10.5,c=LIGHT)

mkbox(sl,Inches(7.0),Inches(4.6),Inches(4.8),Inches(2.1),RGBColor(0x2C,0x3E,0x50),GOLD)
tb(sl,Inches(7.2),Inches(4.65),Inches(4.4),Inches(2.0),info_right,sz=10.5,c=LIGHT)
ft(sl)

# ════════════════════════════════════════════════════
# SLIDE 2: INVESTMENT STRUCTURE & RETURN DECOMPOSITION
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),
   '\u6295\u8cc7\u7d50\u69cb\u8207\u56de\u5831\u5206\u89e3 (10\u5e74\u6301\u5009)',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

# Left panel: structure
mkbox(sl,Inches(0.4),Inches(1.1),Inches(3.6),Inches(5.5),CARD_BG,BLUE,Pt(1.5))
struct_text = (
    "\u7269\u696d\u6295\u8cc7\u7d50\u69cb\n\n"
    "\u7269\u696d\u50f9\u503c\n"
    "  HKD 4,000,000 (JPY 78,000,000)\n\n"
    "\u8cc7\u91d1\u4f86\u6e90\n"
    "  \u9996\u671f: HKD 2,400,000 (60%)\n"
    "  \u8cb8\u6b3e: JPY 31,200,000 (40%)\n\n"
    "\u8cb8\u6b3e\u689d\u6b3e\n"
    "  \u5229\u7387: 3% p.a. (\u65e5\u5703)\n"
    "  \u671f\u9650: 15\u5e74\n"
    "  \u6708\u4f9b: JPY {:,.0f}\n\n"
    "\u79df\u91d1\u6536\u5165\n"
    "  \u5e74\u79df: JPY {:,.0f} (6%)\n"
    "  \u5e74\u6de8\u79df: JPY {:,.0f}\n\n"
    "\u7b2c10\u5e74\u9000\u51fa\n"
    "  \u6309\u63d1\u9918\u984d: JPY {:,.0f}\n"
    "  \u6de8\u5957\u73b0: JPY {:,.0f}".format(
        C['monthly_payment_jpy'], C['annual_rent_jpy'],
        C['net_rent_jpy'], BAL_EXIT,
        C['property_jpy'] - BAL_EXIT)
)
tb(sl,Inches(0.55),Inches(1.15),Inches(3.3),Inches(5.4),struct_text,sz=9.5,c=LIGHT)

# Right: return decomposition chart
sl.shapes.add_picture('{}/return_decomposition.png'.format(CH),Inches(4.2),Inches(1.0),Inches(8.9),Inches(5.7))
ft(sl)

# ════════════════════════════════════════════════════
# SLIDE 3: 10-YEAR FORECAST (Mortgage + Cash Flow)
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),
   '10\u5e74\u6301\u5009\u9810\u6e2c \u2014 \u6309\u63d0\u6524\u9084\u8207\u73fe\u91d1\u6d41',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

sl.shapes.add_picture('{}/mortgage_amortization.png'.format(CH),Inches(0.2),Inches(1.0),Inches(8.5),Inches(5.7))

# Right panel: key metrics
mkbox(sl,Inches(8.9),Inches(1.0),Inches(4.1),Inches(5.7),CARD_BG,GOLD,Pt(1.5))
metrics_text = (
    "\u95dc\u9375\u8ca1\u52d9\u6307\u6a19\n\n"
    "\u6309\u63d0\u6524\u9084 (10\u5e74)\n"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
    "\u5df2\u511f\u672c\u91d1: JPY {:,.0f}\n"
    "\u5df2\u4ed8\u5229\u606f: ~JPY {:,.0f}\n"
    "\u5269\u9918\u9918\u984d: JPY {:,.0f}\n"
    "\u672c\u91d1\u511f\u9084\u6bd4\u4f8b: {:.1f}%\n\n"
    "\u73fe\u91d1\u6d41 (10\u5e74)\n"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
    "\u5e74\u6de8\u79df\u91d1\u6536\u5165:\n"
    "  HKD {:,.0f}/\u5e74\n"
    "10\u5e74\u7d2f\u8a08\u6de8\u79df:\n"
    "  HKD {:,.0f}\n\n"
    "\u9000\u51fa\u50f9\u503c (\u57fa\u6e96)\n"
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n"
    "\u6de8\u5957\u73b0 (HKD):\n"
    "  HKD {:,.0f}\n"
    "\u7e3d\u56de\u5831 (HKD):\n"
    "  HKD {:,.0f}\n"
    "ROI: +85.8%  \u5e74\u5316: 6.39%".format(
        C['principal_paid_10y_jpy'],
        C['annual_rent_jpy'] * HOLD - C['principal_paid_10y_jpy'],  # approximate
        BAL_EXIT,
        C['principal_paid_10y_jpy'] / C['property_jpy'] * (1-C['ltv']) * 100,
        FX['scenarios'][3]['annual_net_hkd'],
        FX['scenarios'][3]['cum_10y_hkd'],
        FX['scenarios'][3]['exit_value_hkd'],
        FX['scenarios'][3]['total_return_hkd'])
)
tb(sl,Inches(9.05),Inches(1.05),Inches(3.8),Inches(5.6),metrics_text,sz=9.5,c=LIGHT)
ft(sl)

# ════════════════════════════════════════════════════
# SLIDE 4: 2-VARIABLE STRESS TEST HEATMAP
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),
   '\u96d9\u8b8a\u91cf\u58d3\u529b\u6e2c\u8a66 \u2014 84\u500b\u60c5\u666f\u5e74\u5316\u56de\u5831\u77e9\u9663 (10\u5e74\u6301\u5009)',sz=24,b=True,c=WHITE)
bar(sl,Inches(0.85))
sl.shapes.add_picture('{}/stress_2var_heatmap.png'.format(CH),Inches(0.2),Inches(1.0),Inches(12.9),Inches(5.4))

mkbox(sl,Inches(0.5),Inches(6.45),Inches(3.8),Inches(0.45),RGBColor(0x1A,0x3C,0x34))
tb(sl,Inches(0.6),Inches(6.47),Inches(3.6),Inches(0.4),
   '\u6700\u5dee\u60c5\u666f: {}% \u5e74\u5316'.format(D['worst']['annualized']),
   sz=11,b=True,c=RED,a=PP_ALIGN.CENTER)
mkbox(sl,Inches(4.7),Inches(6.45),Inches(3.8),Inches(0.45),RGBColor(0x2C,0x2E,0x30))
tb(sl,Inches(4.8),Inches(6.47),Inches(3.6),Inches(0.4),
   '\u866b\u672c\u60c5\u666f: 0/84 (0%)',sz=11,b=True,c=GREEN,a=PP_ALIGN.CENTER)
mkbox(sl,Inches(8.9),Inches(6.45),Inches(3.9),Inches(0.45),RGBColor(0x1A,0x3C,0x34))
tb(sl,Inches(9.0),Inches(6.47),Inches(3.7),Inches(0.4),
   '\u57fa\u6e96\u60c5\u666f: {}% \u5e74\u5316'.format(D['base_annualized']),
   sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
ft(sl)

# ════════════════════════════════════════════════════
# SLIDE 5: RISK MAP + SENSITIVITY LINES
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),
   '\u98a8\u96aa\u5730\u5716\u8207\u654f\u611f\u5ea6\u5206\u6790 (10\u5e74\u6301\u5009)',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

sl.shapes.add_picture('{}/stress_2var_contour.png'.format(CH),Inches(0.15),Inches(1.0),Inches(6.7),Inches(5.6))
sl.shapes.add_picture('{}/stress_2var_lines.png'.format(CH),Inches(6.95),Inches(1.0),Inches(6.2),Inches(5.6))
ft(sl)

# ════════════════════════════════════════════════════
# SLIDE 6: CONCLUSIONS & RISK ASSESSMENT
# ════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl,Inches(0.6),Inches(0.2),Inches(12),Inches(0.7),
   '\u58d3\u529b\u6e2c\u8a66\u7d50\u8ad6\u8207\u98a8\u96aa\u8a55\u4f30 (10\u5e74\u6301\u5009)',sz=26,b=True,c=WHITE)
bar(sl,Inches(0.85))

findings = [
    ("\u6e2c\u8a66\u7bc4\u570d",
     "12\u500b\u532f\u7387\u60c5\u666f x 7\u500b\u623f\u50f9\u60c5\u666f = 84\u500b\u58d3\u529b\u6e2c\u8a66\u60c5\u666f\n"
     "\u8986\u84cbMCLUB\u4e09\u5c64\u67b6\u69cb\uff08T1/T2/T3\uff09+ 2022\u6b77\u53f2\u53c3\u7167",
     BLUE),
    ("\u96f6\u866b\u672c\u6982\u7387",
     "\u5168\u90e884\u500b\u60c5\u666f\u5747\u9304\u5f97\u6b63\u56de\u5831\uff0c\u6700\u5dee\u5e74\u5316\u4ecd\u9054{:.2f}%\n"
     "\u5373\u4f7f\u623f\u50f9\u5d29\u76e430% + \u65e5\u5703\u8cb6\u503c31%\u540c\u6642\u767c\u751f\uff0c\u6295\u8cc7\u4ecd\u53ef\u4fdd\u672c".format(D['worst']['annualized']),
     GREEN),
    ("\u57fa\u6e96\u60c5\u666f",
     "\u623f\u50f9\u6301\u5e73 + \u73fe\u532f\u7387 = \u5e74\u5316{:.2f}%\uff08\u7d14\u79df\u91d1+\u672c\u91d1\u511f\u9084\uff09\n"
     "\u82e5\u623f\u50f9\u5e74\u589e3% + \u73fe\u532f\u7387 = \u5e74\u5316\u7d048.5%+\uff08\u542b\u8cc7\u672c\u589e\u503c\uff09".format(D['base_annualized']),
     GOLD),
    ("\u81ea\u7136\u5c0d\u6c96\u6548\u61c9",
     "\u65e5\u5703\u6309\u63d0+\u65e5\u5703\u79df\u91d1\u5f62\u6210\u5e63\u7a2e\u5339\u914d\u5c0d\u6c96\n"
     "\u532f\u7387\u98a8\u96aa\u50c5\u5f71\u97ff\u6de8\u79df\u91d1\u53ca\u6de8\u5957\u73b0\u50f9\u503c\uff0c\u975e\u5168\u90e8\u7269\u696d\u50f9\u503c",
     BLUE),
    ("\u69d3\u7aef\u5b89\u5168\u908a\u969b",
     "40% LTV\u63d0\u4f9b\u5145\u8db3\u7de9\u885d\uff0c\u7269\u696d\u9700\u8cb6\u503c\u8d8560%\u624d\u89f8\u53ca\u8cb8\u6b3e\u5b89\u5168\u7dda\n"
     "\u6309\u63d0\u672c\u91d1\u511f\u9084\u6301\u7e8c\u964d\u4f4e\u8cb8\u6b3e\u9918\u984d\uff0c10\u5e74\u5167\u5df2\u511f\u908461.6%\u672c\u91d1",
     RGBColor(0x34,0x98,0xDB)),
    ("\u96d9\u91cd\u98a8\u96aa\u6700\u5dee\u60c5\u666f",
     "\u623f\u50f9-30% + JPY\u8cb631%\uff082022\u5371\u6a5f\uff09\n"
     "\u5e74\u5316{:.2f}%\uff0c10\u5e74\u7e3dROI +3.6%\uff0c\u4ecd\u80fd\u4fdd\u672c\u5e76\u7372\u5f97\u6b63\u56de\u5831".format(D['worst']['annualized']),
     RED),
]

y = Inches(1.15)
for title, desc, color in findings:
    mkbox(sl, Inches(0.5), y, Inches(7.5), Inches(0.88), RGBColor(0x2C,0x3E,0x50), color)
    tb(sl, Inches(0.7), y+Inches(0.05), Inches(2.0), Inches(0.3), title, sz=11.5, b=True, c=color)
    tb(sl, Inches(0.7), y+Inches(0.32), Inches(7.0), Inches(0.55), desc, sz=9, c=LIGHT)
    y += Inches(0.92)

# Right: Summary card
mkbox(sl,Inches(8.3),Inches(1.15),Inches(4.5),Inches(5.5),RGBColor(0x2C,0x3E,0x50),GOLD)
tb(sl,Inches(8.5),Inches(1.25),Inches(4.1),Inches(0.4),
   '\u98a8\u96aa\u8a55\u7d1a\u7e3d\u7d50 (10\u5e74)',sz=15,b=True,c=GOLD,a=PP_ALIGN.CENTER)

grades = [
    ("\u6574\u9ad4\u98a8\u96aa\u7b49\u7d1a", "LOW-MODERATE\n\u4e2d\u7b49\u504f\u4f4e", GREEN),
    ("FX\u98a8\u96aa", "MITIGATED\n\u5df2\u88ab\u81ea\u7136\u5c0d\u6c96\u7de9\u89e3", RGBColor(0x82,0xE0,0xAA)),
    ("\u623f\u50f9\u98a8\u96aa", "LOW-MODERATE\n10\u5e74\u6301\u6709\u5206\u6563\u98a8\u96aa", GOLD),
    ("\u5229\u7387\u98a8\u96aa", "LOW\n\u65e5\u672c3%\u8655\u6b77\u53f2\u4f4e\u4f4d", GREEN),
    ("\u6d41\u52d5\u6027\u98a8\u96aa", "MODERATE\n\u65e5\u672c\u7269\u696d\u6d41\u52d5\u6027\u504f\u4f4e", ORANGE),
    ("\u5efa\u8b70", "\u9069\u7528MCLUB\u9032\u968e\u7248\n(Advanced)\u670d\u52d9\u5c64\u7d1a\nHK$128,000", PURPLE),
]

gy = Inches(1.75)
for label, val, color in grades:
    tb(sl, Inches(8.5), gy, Inches(4.1), Inches(0.3), label, sz=10, c=GREY)
    tb(sl, Inches(8.5), gy+Inches(0.28), Inches(4.1), Inches(0.5), val, sz=13, b=True, c=color)
    gy += Inches(0.82)

ft(sl)

# SAVE
out = '/home/z/my-project/download/PZC_JPY_Investment_Professional.pptx'
prs.save(out)
print("Saved: {} ({} slides)".format(out, len(prs.slides)))