# -*- coding: utf-8 -*-
"""
PZC Group — JPY Property FX Risk Scenario Analysis PPTX
8 slides with embedded charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json, os

# ── Load data ──
with open('/home/z/my-project/download/ppt_charts/fx_analysis_data.json', 'r', encoding='utf-8') as f:
    data = json.load(f)

s = data['summary']
scenarios = data['scenarios']
CHART_DIR = '/home/z/my-project/download/ppt_charts'

# ── Color constants ──
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
BLUE   = RGBColor(0x1E, 0x3A, 0x5F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF2, 0xE8)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
GREEN  = RGBColor(0x2E, 0xCC, 0x71)
GREY   = RGBColor(0x95, 0xA5, 0xA6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLD_W = Inches(13.333)
SLD_H = Inches(7.5)

# ── Helpers ──
def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Sarasa Mono SC'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_gold_bar(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLD_W, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_footer(slide, text="PZC Group 百盛大通 — Family Office | JPY Property FX Risk Analysis"):
    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                 text, font_size=9, color=GREY, alignment=PP_ALIGN.CENTER)

def fmt_num(n, prefix='', suffix=''):
    if abs(n) >= 1e6:
        return f"{prefix}{n/1e6:,.2f}M{suffix}"
    elif abs(n) >= 1e3:
        return f"{prefix}{n/1e3:,.1f}K{suffix}"
    else:
        return f"{prefix}{n:,.0f}{suffix}"

# ══════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_gold_bar(slide, Inches(2.8))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             'PZC GROUP 百盛大通', font_size=18, color=GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.5),
             '日本物業投資 — 匯率風險情景分析', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.8),
             'JPY Property Investment — FX Risk Scenario Analysis', font_size=18, color=GREY, alignment=PP_ALIGN.CENTER)

# Key params on cover
params_text = (
    f"物業價值: HKD {s['property_hkd']:,} (JPY {s['property_jpy']:,.0f})  |  "
    f"按揭: 40% LTV @ 3%  |  期限: 15年  |  "
    f"年回報率: 6%  |  基準匯率: {s['fx_rate']} JPY/HKD"
)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             params_text, font_size=13, color=GOLD, alignment=PP_ALIGN.CENTER)

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 2 — INVESTMENT OVERVIEW (key metrics cards)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '投資架構概覽', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

# 4 key metric cards in a row
cards = [
    ('物業價值', f"HKD {s['property_hkd']:,}", f"JPY {s['property_jpy']:,.0f}", BLUE),
    ('首期投入', f"HKD {s['down_payment_hkd']:,.0f}", '60% 自有資金', RGBColor(0x2C, 0x3E, 0x50)),
    ('按揭貸款', f"JPY {s['loan_jpy']:,.0f}", '40% LTV @ 3% p.a.', RGBColor(0x8E, 0x44, 0xAD)),
    ('按揭月供', f"JPY {s['monthly_payment_jpy']:,.0f}", f"年供 JPY {s['annual_payment_jpy']:,.0f}", RGBColor(0xD3, 0x54, 0x00)),
]

card_w = Inches(2.8)
card_h = Inches(1.8)
start_x = Inches(0.5)
gap = Inches(0.3)

for i, (title, val, sub, color) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    shape = add_shape(slide, x, Inches(1.4), card_w, card_h, color, GOLD, Pt(2))
    add_text_box(slide, x + Inches(0.2), Inches(1.5), card_w - Inches(0.4), Inches(0.5),
                 title, font_size=13, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(0.6),
                 val, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), card_w - Inches(0.4), Inches(0.5),
                 sub, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Income breakdown
y_start = Inches(3.6)
add_text_box(slide, Inches(0.6), y_start, Inches(6), Inches(0.5),
             '收入結構分析', font_size=18, bold=True, color=WHITE)

income_items = [
    ('年租金收入 (毛)', 'JPY {:,.0f}'.format(s['gross_rent_jpy']), '物業價值 x 6%'),
    ('年按揭支出', 'JPY {:,.0f}'.format(s['annual_payment_jpy']), 'JPY {:,.0f}/月 x 12'.format(s['monthly_payment_jpy'])),
    ("年淨租金收入", "JPY {:,.0f}".format(s["net_rent_jpy"]), "港元: HKD {:,.0f}/年".format(scenarios[3]["annual_net_hkd"])),
    ("淨租金回報率", "{:.2f}%".format(scenarios[3]["net_yield_pct"]), "相對首期投入 (基準匯率)"),
]

for j, (label, val, note) in enumerate(income_items):
    y = y_start + Inches(0.5) + j * Inches(0.6)
    add_text_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.4), label, font_size=12, color=LIGHT)
    add_text_box(slide, Inches(3.6), y, Inches(2), Inches(0.4), val, font_size=13, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(5.8), y, Inches(2.5), Inches(0.4), note, font_size=10, color=GREY)

# Amortization chart on right
slide.shapes.add_picture(f'{CHART_DIR}/mortgage_amortization.png',
                         Inches(6.8), Inches(3.3), Inches(6.2), Inches(3.4))

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 3 — FX SCENARIO TABLE
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '匯率情景分析總覽', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

add_text_box(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
             '假設物業日圓價值不變，測試不同JPY/HKD匯率對港元回報的影響',
             font_size=12, color=GREY)

# Table
rows = len(scenarios) + 1
cols = 7
tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5)).table

# Column widths
widths = [Inches(2.0), Inches(1.3), Inches(1.6), Inches(1.8), Inches(1.6), Inches(2.0), Inches(2.0)]
for i, w in enumerate(widths):
    tbl.columns[i].width = w

headers = ['匯率情景', 'JPY/HKD\n匯率', '每年淨收入\n(HKD)', '淨回報率\n(年租)', '15年累計\n租金(HKD)', '物業套現\n(HKD)', '15年總\nROI']
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.font.name = 'Sarasa Mono SC'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

row_colors = [
    RGBColor(0x1A, 0x3C, 0x34),  # dark green
    RGBColor(0x1A, 0x35, 0x2E),
    RGBColor(0x25, 0x35, 0x30),
    RGBColor(0x2C, 0x2E, 0x30),  # neutral
    RGBColor(0x35, 0x2E, 0x25),
    RGBColor(0x3A, 0x25, 0x20),
    RGBColor(0x3C, 0x1A, 0x1A),  # dark red
]

for r, sc in enumerate(scenarios):
    row_data = [
        sc['label'],
        f"{sc['fx_rate']:.2f}",
        f"HKD {sc['annual_net_hkd']:,.0f}",
        f"{sc['net_yield_pct']:.2f}%",
        f"HKD {sc['cum_15y_hkd']:,.0f}",
        f"HKD {sc['exit_value_hkd']:,.0f}",
        f"{sc['roi_pct']:+.1f}%",
    ]
    for c, val in enumerate(row_data):
        cell = tbl.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.name = 'Sarasa Mono SC'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_colors[r]
    # Highlight ROI column color
    roi_cell = tbl.cell(r+1, 6)
    for p in roi_cell.text_frame.paragraphs:
        if sc['roi_pct'] > 130:
            p.font.color.rgb = GREEN
        elif sc['roi_pct'] < 100:
            p.font.color.rgb = RED
        else:
            p.font.color.rgb = GOLD

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 4 — Annual Net Income Chart
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '每年淨租金收入 — 匯率情景對比', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

slide.shapes.add_picture(f'{CHART_DIR}/fx_annual_income.png',
                         Inches(0.5), Inches(1.3), Inches(7.5), Inches(4.0))

# Key insight box
box = add_shape(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(5.2),
                RGBColor(0x2C, 0x3E, 0x50), GOLD, Pt(2))

insights = [
    ("核心發現", True, 16, GOLD),
    ("", False, 8, WHITE),
    ("JPY升值情景:", True, 13, GREEN),
    ("日圓每升值10%，淨租金收入增加約HKD 12K/年，來自日圓兌換更多港元。", False, 11, WHITE),
    ("", False, 6, WHITE),
    ("JPY貶值情景:", True, 13, RED),
    ("日圓每貶值10%，淨租金收入減少約HKD 10K/年，日圓購買力下降。", False, 11, WHITE),
    ("", False, 6, WHITE),
    ("基準情景:", True, 13, GOLD),
    (f"每年淨收入 HKD {scenarios[3]['annual_net_hkd']:,.0f}，淨回報率 {scenarios[3]['net_yield_pct']:.2f}%", False, 11, WHITE),
    ("", False, 6, WHITE),
    ("內在對沖效應:", True, 13, WHITE),
    ("按揭以日圓計算、租金以日圓收取，形成自然對沖。匯率風險主要影響淨租金收入及退出價值。", False, 11, LIGHT),
]

y_pos = Inches(1.5)
for text, bold, size, color in insights:
    if text:
        add_text_box(slide, Inches(8.6), y_pos, Inches(4.0), Inches(0.5),
                     text, font_size=size, bold=bold, color=color)
    y_pos += Inches(0.35) if size > 10 else Inches(0.15)

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 5 — 15-Year Total Return
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '15年投資總回報 — 不同匯率情景', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

slide.shapes.add_picture(f'{CHART_DIR}/fx_total_return.png',
                         Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.3))

add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             '* 總回報 = 15年累計淨租金 + 物業套現價值（假設日圓物業價格不變）  |  投入本金: HKD 2,400,000',
             font_size=10, color=GREY, alignment=PP_ALIGN.CENTER)

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 6 — FX Sensitivity
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '匯率敏感度分析', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

slide.shapes.add_picture(f'{CHART_DIR}/fx_sensitivity.png',
                         Inches(0.3), Inches(1.3), Inches(8.0), Inches(4.3))

# Risk summary cards on right
risk_cards = [
    ('最佳情景', 'JPY升值+20%', f"年化回報: {scenarios[0]['annualized_roi']:.1f}%", GREEN),
    ('基準情景', '現匯率', f"年化回報: {scenarios[3]['annualized_roi']:.1f}%", GOLD),
    ('最差情景', 'JPY貶值-20%', f"年化回報: {scenarios[6]['annualized_roi']:.1f}%", RED),
]

for i, (title, sub, metric, color) in enumerate(risk_cards):
    y = Inches(1.4) + i * Inches(1.8)
    shape = add_shape(slide, Inches(8.6), y, Inches(4.2), Inches(1.5), RGBColor(0x2C, 0x3E, 0x50), color, Pt(2))
    add_text_box(slide, Inches(8.8), y + Inches(0.1), Inches(3.8), Inches(0.4),
                 title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.8), y + Inches(0.5), Inches(3.8), Inches(0.4),
                 sub, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.8), y + Inches(0.9), Inches(3.8), Inches(0.4),
                 metric, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key note
add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
             '關鍵: 即使在最差情景（日圓貶值20%），年化回報仍達4.55%，顯示按揭槓杆在貶值環境中提供了一定緩衝。',
             font_size=12, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 7 — Summary Dashboard
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '關鍵指標摘要', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

slide.shapes.add_picture(f'{CHART_DIR}/fx_summary_dashboard.png',
                         Inches(0.3), Inches(1.3), Inches(12.6), Inches(5.3))

add_footer(slide)

# ══════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSIONS & RECOMMENDATIONS
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             '結論與建議', font_size=28, bold=True, color=WHITE)
add_gold_bar(slide, Inches(1.0))

# Left column - findings
findings = [
    ("投資回報概要", [
        f"基準情景下15年總回報率達 {scenarios[3]['roi_pct']:+.1f}%，年化 {scenarios[3]['annualized_roi']:.1f}%",
        f"最佳情景（JPY+20%）年化回報 {scenarios[0]['annualized_roi']:.1f}%",
        f"最差情景（JPY-20%）年化回報仍達 {scenarios[6]['annualized_roi']:.1f}%",
        f"按揭槓杆（40% LTV）有效提升資金運用效率",
    ]),
    ("匯率風險特徵", [
        "日圓按揭 + 日圓租金形成自然對沖，大幅降低匯率波動影響",
        "匯率風險主要體現在淨租金收入（約佔總回報25-30%）及退出價值",
        "日圓升值對投資者有利，貶值則壓縮回報但仍維持正回報",
    ]),
]

y = Inches(1.3)
for title, items in findings:
    add_text_box(slide, Inches(0.6), y, Inches(6), Inches(0.4),
                 title, font_size=16, bold=True, color=GOLD)
    y += Inches(0.45)
    for item in items:
        add_text_box(slide, Inches(0.9), y, Inches(5.7), Inches(0.35),
                     f"• {item}", font_size=11, color=WHITE)
        y += Inches(0.35)
    y += Inches(0.2)

# Right column - recommendations
box = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.3),
                RGBColor(0x2C, 0x3E, 0x50), GOLD, Pt(2))

add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.2), Inches(0.5),
             '風險管理建議', font_size=16, bold=True, color=GOLD)

recs = [
    ("1. 自然對沖優勢", "日圓按揭與日圓租金的幣種匹配已提供內建對沖，建議維持此結構，不需額外進行複雜的遠期合約對沖。"),
    ("2. 分階段匯率轉換", "建議將每年的淨租金收入分4次（季度）轉換為港元，降低單次匯率波動的影響。"),
    ("3. 退出時機規劃", "15年後物業出售時匯率不確定性最大，建議設定匯率目標區間（如JPY/HKD 17-22），在區間內分批套現。"),
    ("4. 利率風險關注", "目前日本按揭利率3%處歷史低位，若未來日銀加息，需預留利率上升1-2%的緩衝空間。"),
    ("5. 定期檢視", "建議每半年檢視一次匯率走勢及租金回報，必要時調整策略。"),
]

ry = Inches(2.0)
for title, desc in recs:
    add_text_box(slide, Inches(7.3), ry, Inches(5.2), Inches(0.35),
                 title, font_size=12, bold=True, color=WHITE)
    ry += Inches(0.35)
    add_text_box(slide, Inches(7.3), ry, Inches(5.2), Inches(0.5),
                 desc, font_size=10, color=LIGHT)
    ry += Inches(0.55)

add_footer(slide)

# ── Save ──
out_path = '/home/z/my-project/download/PZC_FX_Risk_Scenario_JPY_Property.pptx'
prs.save(out_path)
print(f"✓ PPTX saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
